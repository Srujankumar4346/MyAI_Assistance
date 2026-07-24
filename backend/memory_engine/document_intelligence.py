"""
Phase 3 — Document Intelligence

Extracts text from uploaded documents (PDF, DOCX, PPTX, XLSX, TXT),
summarizes it, and stores important facts as enhanced memories.

Supported formats (falls back gracefully if library not installed):
  - PDF → PyPDF2 or pdfminer.six
  - DOCX → python-docx
  - PPTX → python-pptx
  - XLSX → openpyxl
  - TXT → plain text read

After extraction:
  1. Text is chunked into ~500 token segments
  2. Each chunk is scored for importance
  3. High-importance chunks become EnhancedMemory records
  4. Entities are sent to the KnowledgeGraph
  5. A DocumentMemory record is saved with status tracking
"""
import uuid
import io
from typing import List, Optional, Tuple
from datetime import datetime

from backend.database.connection import SessionLocal
from backend.memory_engine.models import DocumentMemory
from backend.memory_engine.neural_memory import neural_memory
from backend.memory_engine.knowledge_graph import knowledge_graph
from backend.utils.logger import logger


# ── Text Extraction ────────────────────────────────────────────────────────────

def _extract_pdf(content: bytes) -> Tuple[str, int]:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(content))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(pages), len(pages)
    except ImportError:
        pass
    try:
        from pdfminer.high_level import extract_text_to_fp
        from pdfminer.layout import LAParams
        output = io.StringIO()
        extract_text_to_fp(io.BytesIO(content), output, laparams=LAParams())
        text = output.getvalue()
        return text, text.count("\x0c") + 1  # form feeds = page breaks
    except ImportError:
        return "[PDF extraction requires PyPDF2 or pdfminer.six]", 0


def _extract_docx(content: bytes) -> Tuple[str, int]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paragraphs), len(doc.paragraphs)
    except ImportError:
        return "[DOCX extraction requires python-docx]", 0


def _extract_pptx(content: bytes) -> Tuple[str, int]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for slide in prs.slides:
            texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            slides_text.append(" ".join(texts))
        return "\n".join(slides_text), len(prs.slides)
    except ImportError:
        return "[PPTX extraction requires python-pptx]", 0


def _extract_xlsx(content: bytes) -> Tuple[str, int]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        rows_all = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(c) for c in row if c is not None)
                if row_text.strip():
                    rows_all.append(row_text)
        return "\n".join(rows_all), len(wb.worksheets)
    except ImportError:
        return "[XLSX extraction requires openpyxl]", 0


def extract_text(filename: str, content: bytes) -> Tuple[str, str, int]:
    """Returns (text, file_type, page_count)."""
    fn_lower = filename.lower()
    if fn_lower.endswith(".pdf"):
        text, pages = _extract_pdf(content)
        return text, "pdf", pages
    elif fn_lower.endswith(".docx"):
        text, pages = _extract_docx(content)
        return text, "docx", pages
    elif fn_lower.endswith(".pptx"):
        text, pages = _extract_pptx(content)
        return text, "pptx", pages
    elif fn_lower.endswith((".xlsx", ".xls")):
        text, pages = _extract_xlsx(content)
        return text, "xlsx", pages
    else:
        # Plain text fallback
        text = content.decode("utf-8", errors="replace")
        return text, "txt", 1


# ── Text Chunking ──────────────────────────────────────────────────────────────

def chunk_text(text: str, chunk_size: int = 500) -> List[str]:
    """Split text into ~chunk_size word chunks at sentence boundaries."""
    import re
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks: List[str] = []
    current: List[str] = []
    word_count = 0
    for sent in sentences:
        words = len(sent.split())
        if word_count + words > chunk_size and current:
            chunks.append(" ".join(current))
            current = []
            word_count = 0
        current.append(sent)
        word_count += words
    if current:
        chunks.append(" ".join(current))
    return [c for c in chunks if len(c.strip()) > 50]


# ── Document Intelligence Engine ───────────────────────────────────────────────

class DocumentIntelligence:

    async def process_document(
        self,
        filename: str,
        content: bytes,
        user_id: int,
        auto_extract_memories: bool = True,
    ) -> dict:
        """
        Full pipeline: extract → chunk → score → store memories → knowledge graph.
        Returns summary of what was done.
        """
        doc_id = str(uuid.uuid4())
        db = SessionLocal()
        try:
            # ── Extract text ──────────────────────────────────────────────────
            text, file_type, page_count = extract_text(filename, content)
            text = text.strip()

            # ── Save DocumentMemory record ────────────────────────────────────
            doc_mem = DocumentMemory(
                id=doc_id,
                user_id=user_id,
                filename=filename,
                file_type=file_type,
                file_size_bytes=len(content),
                page_count=page_count,
                extracted_text=text[:50000],  # truncate for storage
                status="processing",
            )
            db.add(doc_mem)
            db.commit()

            if not auto_extract_memories or not text:
                doc_mem.status = "done"
                db.commit()
                return {"doc_id": doc_id, "memories_extracted": 0, "file_type": file_type}

            # ── Chunk and store memories ──────────────────────────────────────
            chunks = chunk_text(text, chunk_size=300)
            memories_created = 0

            for chunk in chunks[:50]:  # max 50 memories per doc
                result = await neural_memory.store_memory(
                    content=chunk,
                    user_id=user_id,
                    memory_type="document",
                    category="documents",
                    source="document",
                    project_name=filename,
                )
                if result and result.get("status") == "created":
                    memories_created += 1

            # ── Extract entities to knowledge graph ───────────────────────────
            await knowledge_graph.ingest_text(text[:5000], user_id)

            # ── Update document record ────────────────────────────────────────
            doc_mem.memories_extracted = memories_created
            doc_mem.status = "done"
            db.commit()

            logger.info(f"[DocIntel] Processed '{filename}': {memories_created} memories created.")
            return {
                "doc_id": doc_id,
                "filename": filename,
                "file_type": file_type,
                "page_count": page_count,
                "memories_extracted": memories_created,
                "text_length": len(text),
            }
        except Exception as e:
            db.rollback()
            logger.error(f"[DocIntel] process_document failed: {e}")
            if 'doc_mem' in dir():
                try:
                    doc_mem.status = "failed"
                    db.commit()
                except Exception:
                    pass
            return {"doc_id": doc_id, "error": str(e), "memories_extracted": 0}
        finally:
            db.close()

    async def list_documents(self, user_id: int) -> List[dict]:
        db = SessionLocal()
        try:
            docs = db.query(DocumentMemory).filter(
                DocumentMemory.user_id == user_id
            ).order_by(DocumentMemory.created_at.desc()).all()
            return [
                {
                    "id": d.id,
                    "filename": d.filename,
                    "file_type": d.file_type,
                    "page_count": d.page_count,
                    "memories_extracted": d.memories_extracted,
                    "status": d.status,
                    "created_at": d.created_at.isoformat(),
                }
                for d in docs
            ]
        finally:
            db.close()


document_intelligence = DocumentIntelligence()
